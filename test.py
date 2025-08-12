from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import date

TITLE = "Algo Trading SaaS Startup Roadmap (Corrected)"
SUBTITLE = f"Plan, metrics, financials, and compliance overview — v1.0 — {date.today().isoformat()}"
OUTPUT_FILE = "SaaS_AlgoTrading_Roadmap_Corrected.pptx"

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line if isinstance(line, str) else line[0]
        p.level = 0
        if isinstance(line, (list, tuple)) and len(line) > 1:
            for sub in line[1:]:
                sp = tf.add_paragraph()
                sp.text = f"- {sub}"
                sp.level = 1
    return slide

def add_table_slide(prs, title, col_headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = title
    cols = len(col_headers)
    rows_count = len(rows) + 1
    left, top, width, height = Inches(0.5), Inches(1.6), Inches(9), Inches(4.5)
    table = slide.shapes.add_table(rows_count, cols, left, top, width, height).table

    # headers
    for j, h in enumerate(col_headers):
        cell = table.cell(0, j)
        cell.text = h

    # rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)
    return slide

def main():
    prs = Presentation()

    # 1. Title
    add_title_slide(prs, TITLE, SUBTITLE)

    # 2. Executive Summary
    add_bullets_slide(prs, "Executive Summary", [
        "Promising direction: SaaS-first with exportable algo bots and signals; prop trading used to validate and showcase.",
        "Critical corrections added: robust validation (walk-forward, PSR/DSR), cost-aware metrics, risk and compliance gates.",
        "Phased expansion: Phase 0 (R&D gate) → Phase 1 (BD + global SaaS) → Phase 2 (India) → Phase 3 (UAE).",
        "Financial model included: break-even, ROI, Year-1 scenarios (conservative/base/extraordinary), conservative and worst-case loss."
    ])

    # 3. Offerings & Revenue
    add_bullets_slide(prs, "Offerings & Revenue Streams", [
        "SaaS: MT4/MT5 exportable algo bots and strategy signals (generic, not personalized advice).",
        "PaaS: NixOS/Kubernetes server setup and stable backtesting/execution environments for traders.",
        "AI utilities: pay-as-you-go analysis tools (hourly/daily/weekly/monthly).",
        "Prop trading: DSE low-frequency, XAUUSD intraday (R&D), global crypto strategic reserve.",
        "Future: enterprise deployments and APIs; partner broker integrations (introducing-broker deals)."
    ])

    # 4. Tech & Data Stack
    add_bullets_slide(prs, "Tech, Data, and MLOps", [
        "Core stack: Python, C++/Rust for latency, Kafka, NumPy/Pandas, Keras/PyTorch, MQL5/MT5, Next.js/Redux, Grafana, Wazuh SIEM.",
        "MLOps: MLflow (experiments), DVC (data/versioning), model registry, feature store, CI/CD & canary, drift detection and rollback.",
        "Observability & Sec: Prometheus/Grafana, Loki/OpenTelemetry, Vault for secrets, SBOM/SCA checks.",
        "Data quality: point-in-time, corporate-actions-adjusted equities with delisted names; exchange-specific crypto data; license compliance."
    ])

    # 5. Model Validation Corrections
    add_bullets_slide(prs, "Model Validation Corrections", [
        "Use walk-forward analysis with purged, embargoed cross-validation; hold out 30–40% out-of-sample; evaluate by regimes (bull/bear/sideways).",
        "All metrics are net-of-costs: commissions, spreads, funding/overnight, slippage, and taxes.",
        "Robustness: Probabilistic/Deflated Sharpe (PSR/DSR), White’s Reality Check, bootstrapped confidence intervals.",
        "Define expectancy properly: win rate plus payoff ratio (avg win/avg loss ≥ 1.2)."
    ])

    # 6. Performance Gates & Tiers
    add_bullets_slide(prs, "Performance Gates & Tiers (Net of Costs)", [
        "Promotion to production requires all:",
        [
            "Sharpe ≥ 1.6 (Pro), ≥ 2.0 (Top tier).",
            "Max drawdown ≤ 15% (public-facing strategies).",
            "PSR ≥ 0.95; Calmar ≥ 0.7.",
            "Stable across ≥ 3 market regimes; capacity-tested; slippage-model validated."
        ],
        "Market-viable minimum: Sharpe ≥ 1.2, Max DD ≤ 20%, positive expectancy."
    ])

    # 7. Risk Management & Execution
    add_bullets_slide(prs, "Risk Management & Execution", [
        "Sizing: volatility-scaling or fractional Kelly with caps; 1–2% per-position risk for DSE; hard stops and time stops.",
        "Guardrails: daily 99% VaR ≤ 2% equity; monthly stop-out if DD > 10–15%; circuit breakers.",
        "Execution: event-driven backtester and OMS; model partial fills and latency; limit/iceberg/PO orders; venue-specific microstructure."
    ])

    # 8. Compliance & Licensing Overview
    add_bullets_slide(prs, "Compliance & Licensing Overview", [
        "Bangladesh (BSEC): signals/bots may be considered advice; safest path is software export with clear disclaimers; confirm with counsel.",
        "India (SEBI): Research Analyst/Investment Adviser needed to sell signals/advice to Indian residents; otherwise geo-restrict or partner.",
        "UAE (DIFC/ADGM/VARA): authorization depends on advisory vs prop and crypto scope; plan ahead for onboarding and reporting.",
        "No custody of client funds; robust T&Cs; market data redistribution rights checked."
    ])

    # 9. Phase 0
    add_bullets_slide(prs, "Phase 0 — Mandatory R&D Gate", [
        "Data: Equities 15–20y with delisted names; Crypto 7–10y; XAUUSD tick/1m with realistic spreads.",
        "Validation: walk-forward + purged k-fold; regime splits; PSR/DSR; net-of-costs; capacity tests.",
        "Gate to Phase 1: Sharpe ≥ 1.6, Max DD ≤ 15%, PSR ≥ 0.95, Calmar ≥ 0.7, avg win/avg loss ≥ 1.2, 3-month paper/live-sim stability."
    ])

    # 10. Phase 1
    add_bullets_slide(prs, "Phase 1 — Bangladesh + Global SaaS", [
        "Legal: operate as software provider; avoid personalized advice; finalize T&Cs, disclaimers, data licenses.",
        "Prop trading (own funds): DSE $5k (LFT), Crypto $10k (staged DCA with regime filter), XAUUSD $500 (R&D-only, strict stop-outs).",
        "Infra & costs: owned hardware + VPS/data; fixed run-rate target ~$2.5k/month initially.",
        "Exit gate to Phase 2: 3 months SaaS break-even MRR and live Sharpe ≥ 1.4, Max DD ≤ 12%, India compliance plan ready."
    ])

    # 11. Pricing & Packaging
    add_bullets_slide(prs, "Pricing & Packaging (Indicative)", [
        "Basic $59: 1 symbol/1 bot, delayed signals, community support.",
        "Pro $99: 5 symbols, real-time, cloud backtesting credits, email support.",
        "Advanced $199: 20 symbols, portfolio bots, API access, priority support.",
        "Enterprise: private deployments, SLAs, custom integrations."
    ])

    # 12. Go-to-Market
    add_bullets_slide(prs, "Go-to-Market (GTM)", [
        "Channels: MT5 Market, website, content-led SEO, broker partnerships, technical whitepapers.",
        "Transparency: read-only live dashboards, net-of-costs tear sheets; no performance guarantees.",
        "Support: ticketing with response SLOs; user education and onboarding guides."
    ])

    # 13. Financial Model — Break-even & ROI
    add_bullets_slide(prs, "Financial Model — Break-even & ROI (Illustrative)", [
        "Assumptions: fixed costs ~$2.5k/month; variable fees ~10% of SaaS revenue; initial cash $40k.",
        "Break-even MRR ≈ Fixed / (1 − variable) = 2,500 / 0.90 ≈ $2,780 → ~28 subs at $99 ARPU.",
        "ROI year-1 depends mainly on SaaS subscriber ramp; trading contributes modestly and is risk-capped."
    ])

    # 14. Year-1 Scenarios Table
    cols = ["Scenario", "Subs (M12)", "ARPU ($)", "MRR (M12) ($)", "Annual Net SaaS ($)", "Fixed Costs ($)", "SaaS Op Profit ($)", "Trading Net ($)", "Total Profit ($)", "ROI (%)", "Break-even (mo)"]
    rows = [
        ["Conservative", "50", "79", "3,950", "42,660", "30,000", "12,660", "+1,500/−100 range", "≈13,160", "≈33", "6–7"],
        ["Base", "150", "99", "14,850", "160,380", "42,000", "118,380", "≈+2,000", "≈120,380", "≈301", "4–5"],
        ["Extraordinary", "500", "119", "59,500", "642,600", "120,000", "522,600", "≈+21,875", "≈544,475", "≈1,361", "2–3"],
    ]
    add_table_slide(prs, "Year-1 Financial Scenarios (Illustrative, Net where noted)", cols, rows)

    # 15. Loss Cases
    add_bullets_slide(prs, "Loss Cases & Runway", [
        "Conservative loss: ~20 subs @ $59 → monthly burn ≈ $1,438 → annual ≈ $17,256; runway on $40k ≈ 22–28 months.",
        "Worst case: compliance stop-sale + refunds/legal $15–20k; Crypto −70% (−$7k), DSE −30% (−$1.5k), XAUUSD −100% (−$0.5k).",
        "Cost controls: pause marketing, reduce VPS/data, narrow product scope, extend runway; keep legal reserves."
    ])

    # 16. Phase 2
    add_bullets_slide(prs, "Phase 2 — India (NIFTY)", [
        "Legal: SEBI RA/IA or partner with registered entity; otherwise geo-restrict Indian residents.",
        "Investments (prop): NIFTY intraday/swing $25k; Crypto reserve staged to $20k; WazirX intraday $500 (top-10 tokens only).",
        "Exit gate to Phase 3: 6 months company net profit ≥ 3× monthly fixed costs; UAE compliance readiness."
    ])

    # 17. Phase 3
    add_bullets_slide(prs, "Phase 3 — UAE (Dubai/ADGM/DIFC/VARA)", [
        "Obtain appropriate license for advisory/crypto as applicable; consider prop-only vs public offerings.",
        "HFT research caution: realistic budget often >> $15–30k when including colo, market data, and specialist talent.",
        "Shift 50% R&D to UAE; expand global coverage (SGX, NASDAQ, ASX, LSE, TSE) post product-market fit."
    ])

    # 18. Roadmap & Gates
    add_bullets_slide(prs, "Roadmap Timeline & Phase Gates", [
        "Gate A (Pre-sale): net Sharpe ≥ 1.6, DD ≤ 15%, PSR ≥ 0.95, 3-month live-sim stability.",
        "Gate B (Post-MVP): ≥ 30 paying subs, churn < 8%/mo, tickets < 0.5 per sub/mo, uptime ≥ 99.5%.",
        "Gate C (Pre-India): 3-month consecutive profitability; counsel sign-off on SEBI exposure.",
        "Gate D (Pre-UAE): 6-month profitability; MRR growth ≥ 10% q/q; compliance budget set aside."
    ])

    # 19. KPIs
    add_bullets_slide(prs, "Core KPIs to Track", [
        "SaaS: MRR, ARPU, CAC, LTV, churn, activation time, support load, uptime.",
        "Trading: net Sharpe, Max DD, PSR/DSR, turnover, slippage vs model, capacity utilization.",
        "Compliance/Trust: incidents, data license compliance, audit trails, customer NPS."
    ])

    # 20. Next Steps
    add_bullets_slide(prs, "Next 30–60 Days — Action Plan", [
        "Replace single backtest with walk-forward + purged k-fold; compute PSR/DSR; publish net-of-costs tear sheet.",
        "Finalize T&Cs, disclaimers; confirm BSEC posture; outline SEBI/VARA paths.",
        "Ship SaaS core: auth, billing, license enforcement for EAs, metering, telemetry; publish pricing.",
        "Launch one flagship strategy with transparent live-sim dashboard; start content-led GTM.",
        "Set break-even target (≈28 subs @ $99); track KPI dashboard weekly."
    ])

    # 21. Disclaimers
    add_bullets_slide(prs, "Disclaimers", [
        "Not investment or legal advice. Past performance is not indicative of future results.",
        "All metrics and scenarios are illustrative; validate with counsel and independent financial review.",
        "We do not manage client funds; software is for research/education unless licensed otherwise."
    ])

    prs.save(OUTPUT_FILE)
    print(f"Written: {OUTPUT_FILE}")

if __name__ == "__main__":
    main()